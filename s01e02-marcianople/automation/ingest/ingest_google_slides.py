#!/usr/bin/env python3
"""Download Google Slides deck and extract slide PNGs, text, and speaker notes."""

from __future__ import annotations

import json
import re
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path

import pymupdf
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[3]
EPISODE = ROOT / "s01e02-marcianople"
OUT = EPISODE / "canon" / "google-slides"
DOWNLOAD = OUT / "_download"
SLIDES = OUT / "slides"

DECK_ID = "1480U-gIaQ1XGXPlE_FMH6Vxk3kmRPPvLVocrdzvE98M"
SOURCE_URL = f"https://docs.google.com/presentation/d/{DECK_ID}/edit?usp=sharing"

# Heuristic shot mapping from slide copy (first match wins).
SHOT_RULES: list[tuple[str, list[str], str]] = [
    ("C01", ["river crossing", "series opening"], "Crossing establish"),
    ("C07", ["abuse 1", "abuse 2", "abuse 4", "abuse 5"], "Mistreatment"),
    ("C10", ["line too far", "elder brother", "dogmeat"], "Brother taken"),
    ("C13", ["invite to banquet"], "Banquet invitation"),
    ("H01", ["treachery planned", "planned treachery", "planning something awful"], "Trap prep"),
    ("H02", ["outside gates", "hordes amassing"], "Banquet entry / gates"),
    ("H04", ["chaos at the banquet", "screenprint"], "Banquet chaos (check register)"),
    ("H05", ["fight breaks out", "romans are killed"], "Violence / corridor"),
    ("H03", ["fritigern hostage"], "Hostage"),
    ("H06", ["uniform swap", "weapon heist"], "Uniform + heist"),
    ("H09", ["escaping with the weapons"], "Extraction"),
    ("H10", ["riders return", "weapons arrive"], "Indigo return"),
    ("B01", ["weapons arrive to camp"], "Weapons reach camp"),
    ("B03", ["suit up", "prep for battle", "goth's prepping", "goths' prepping"], "Cloaks fall / prep"),
    ("B05", ["romans marching", "easy slaugher", "easy slaughter"], "Roman approach"),
    ("B02", ["women take infirm", "edge of the forest"], "Families concealed"),
    ("B07", ["gothic horde", "battle scene"], "Battle shock"),
    ("B14", ["standoff", "catches up with fritigern", "battling"], "Fritigern vs Lupicinus"),
    ("B11", ["alaric looks on"], "Alaric watches"),
    ("B15", ["mercy", "yielding", "runs away"], "Mercy"),
    ("B17", ["boy awakens", "preview"], "Inheritance / coda"),
]


def slug(text: str, n: int = 48) -> str:
    s = re.sub(r"[^a-z0-9]+", "-", (text or "slide").lower()).strip("-")
    return (s[:n] or "slide").strip("-")


def suggest_shot(text: str) -> tuple[str | None, str | None]:
    low = text.lower()
    for shot_id, keys, note in SHOT_RULES:
        for key in keys:
            if key in low:
                return shot_id, note
    return None, None


def download_deck() -> tuple[Path, Path, Path]:
    DOWNLOAD.mkdir(parents=True, exist_ok=True)
    base = f"https://docs.google.com/presentation/d/{DECK_ID}/export"
    paths = {
        "pptx": DOWNLOAD / "deck.pptx",
        "pdf": DOWNLOAD / "deck.pdf",
        "txt": DOWNLOAD / "deck.txt",
    }
    for fmt, dest in [("pptx", paths["pptx"]), ("pdf", paths["pdf"]), ("txt", paths["txt"])]:
        url = f"{base}/{fmt}" if fmt != "txt" else f"{base}?format=txt"
        if dest.exists() and dest.stat().st_size > 1000:
            continue
        subprocess.run(["curl", "-sL", "-o", str(dest), url], check=True)
    return paths["pptx"], paths["pdf"], paths["txt"]


def shape_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if t:
            texts.append(t)
    return texts


def extract_pptx(pptx: Path) -> list[dict]:
    prs = Presentation(str(pptx))
    rows: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = shape_texts(slide)
        title = texts[0] if texts else ""
        body = "\n".join(texts[1:]) if len(texts) > 1 else ""
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        combined = "\n".join(filter(None, [title, body, notes]))
        shot_id, reason = suggest_shot(combined)
        rows.append({
            "slide": i,
            "title": title,
            "body": body,
            "notes": notes,
            "suggested_shot_id": shot_id,
            "mapping_reason": reason,
        })
    return rows


def render_pdf_pngs(pdf: Path, out_dir: Path, scale: float = 2.0) -> list[str]:
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = pymupdf.open(str(pdf))
    paths: list[str] = []
    for i in range(doc.page_count):
        page = doc[i]
        pix = page.get_pixmap(matrix=pymupdf.Matrix(scale, scale))
        name = f"slide-{i + 1:02d}.png"
        dest = out_dir / name
        pix.save(str(dest))
        paths.append(dest.relative_to(ROOT).as_posix())
    return paths


def write_review_html(manifest: dict) -> None:
    slides = manifest["slides"]
    cards = []
    for s in slides:
        shot = s.get("suggested_shot_id") or "—"
        cards.append(f"""
    <article class="card">
      <header><strong>Slide {s['slide']:02d}</strong> · shot <code>{shot}</code></header>
      <figure><img src="/{s['image']}" alt="slide {s['slide']}" loading="lazy" /></figure>
      <h2>{s.get('title') or '(no title)'}</h2>
      <pre class="body">{s.get('body') or ''}</pre>
      <details><summary>Speaker notes</summary><pre>{s.get('notes') or '—'}</pre></details>
    </article>""")
    html = f"""<!DOCTYPE html>
<html lang="en"><head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Google Slides — S01E02 storyboard</title>
  <style>
    body {{ font-family: system-ui, sans-serif; background: #111; color: #e8e4dc; margin: 0; padding: 1rem; }}
    .grid {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(320px, 1fr)); gap: 1rem; }}
    .card {{ background: #222; border: 1px solid #333; border-radius: 8px; padding: 0.75rem; }}
    img {{ width: 100%; border-radius: 4px; background: #000; }}
    pre {{ white-space: pre-wrap; font-size: 0.8rem; color: #aaa; }}
    h2 {{ font-size: 1rem; margin: 0.5rem 0; }}
  </style>
</head><body>
  <h1>S01E02 storyboard (Google Slides)</h1>
  <p>Source: <a href="{manifest['source_url']}" style="color:#8af">Google Slides</a> · {manifest['slide_count']} slides</p>
  <div class="grid">{''.join(cards)}</div>
</body></html>"""
    review = EPISODE / "renders" / "reviews" / "google-slides.html"
    review.parent.mkdir(parents=True, exist_ok=True)
    review.write_text(html, encoding="utf-8")


def main() -> int:
    pptx, pdf, txt = download_deck()
    meta = extract_pptx(pptx)
    image_paths = render_pdf_pngs(pdf, SLIDES)

    if len(meta) != len(image_paths):
        print(f"WARN: pptx slides {len(meta)} != pdf pages {len(image_paths)}", file=sys.stderr)

    slides: list[dict] = []
    for i, row in enumerate(meta):
        img = image_paths[i] if i < len(image_paths) else None
        slides.append({**row, "image": img, "slug": slug(row.get("title") or f"slide-{row['slide']}")})

    manifest = {
        "schemaVersion": "1.0.0",
        "deck_id": DECK_ID,
        "source_url": SOURCE_URL,
        "exported_at": datetime.now(timezone.utc).isoformat(),
        "slide_count": len(slides),
        "slides": slides,
        "raw_downloads": {
            "pptx": pptx.relative_to(ROOT).as_posix(),
            "pdf": pdf.relative_to(ROOT).as_posix(),
            "txt": txt.relative_to(ROOT).as_posix(),
        },
    }
    OUT.mkdir(parents=True, exist_ok=True)
    (OUT / "manifest.json").write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    (OUT / "outline.txt").write_text(txt.read_text(encoding="utf-8"), encoding="utf-8")
    write_review_html(manifest)
    print(json.dumps({"slides": len(slides), "out": str(OUT.relative_to(ROOT))}, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
